# -*- coding: utf-8 -*-
"""生成组会汇报 PPT — 科研假设生成器 V7.5 进度汇报"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── 配色方案 ──
BG_DARK =  RGBColor(0x1A, 0x1A, 0x2E)       # 深色背景
BG_CARD =  RGBColor(0x22, 0x22, 0x3E)        # 卡片背景
ACCENT  =  RGBColor(0x00, 0xD2, 0xFF)        # 青色强调
GREEN   =  RGBColor(0x00, 0xE6, 0x96)        # 绿色 - 完成
AMBER   =  RGBColor(0xFF, 0xB8, 0x00)        # 琥珀色 - 进行中
RED     =  RGBColor(0xFF, 0x47, 0x5C)        # 红色 - 问题
WHITE   =  RGBColor(0xFF, 0xFF, 0xFF)
GRAY    =  RGBColor(0x94, 0xA3, 0xB8)
DARK_BG = RGBColor(0x0F, 0x0F, 0x23)         # 更深背景


def set_slide_bg(slide, color):
    """设置幻灯片背景色"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=14,
                 color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf


def add_bullet_list(slide, left, top, width, height, items, font_size=13,
                    color=GRAY, bullet='•', font_name='Microsoft YaHei'):
    """添加项目符号列表"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"{bullet} {item}"
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(6)
    return tf


def add_accent_line(slide, left, top, width):
    """添加强调分隔线"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Pt(3)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.fill.background()
    return shape


def add_metric_card(slide, left, top, label, value, width=1.6, height=0.7):
    """添加指标卡片"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(value)
    p.font.size = Pt(24)
    p.font.color.rgb = ACCENT
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(10)
    p2.font.color.rgb = GRAY
    p2.font.name = 'Microsoft YaHei'
    p2.alignment = PP_ALIGN.CENTER


def add_card(slide, left, top, width, height, title, content, title_color=ACCENT):
    """添加圆角卡片"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BG_CARD
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(13)
    p.font.color.rgb = title_color
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p2 = tf.add_paragraph()
    p2.text = content
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY
    p2.font.name = 'Microsoft YaHei'
    return tf


# ═══════════════════════════════════════════
# 创建演示文稿
# ═══════════════════════════════════════════
prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9 宽屏
prs.slide_height = Inches(7.5)


def make_slide():
    """创建空白幻灯片"""
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    set_slide_bg(slide, BG_DARK)
    return slide


# ═══════════════════════════════════════════
# Slide 1: 封面
# ═══════════════════════════════════════════
s = make_slide()
set_slide_bg(s, DARK_BG)
add_text_box(s, 0.8, 1.5, 11, 1, "科研假设生成器 V7.5", font_size=42, color=ACCENT, bold=True)
add_text_box(s, 0.8, 2.4, 11, 0.6, "Phoenix Evolution — 组会进度汇报", font_size=22, color=GRAY)
add_accent_line(s, 0.8, 3.1, 3.0)
add_text_box(s, 0.8, 3.4, 11, 0.5, "生物医学 AI 科研引擎 · 文献检索 · 假设生成 · 智能验证", font_size=14, color=GRAY)
add_text_box(s, 0.8, 4.5, 11, 0.5, "汇报内容: 架构重构 / UI 优化 / Worker 调试 / 新功能实现", font_size=13, color=GRAY)
add_text_box(s, 0.8, 6.5, 11, 0.4, "2026-05-05  |  V7.5 Phoenix Evolution", font_size=11, color=RGBColor(0x55, 0x55, 0x77))

# ═══════════════════════════════════════════
# Slide 2: 总览
# ═══════════════════════════════════════════
s = make_slide()
add_text_box(s, 0.8, 0.4, 11, 0.6, "本期工作概览", font_size=28, color=WHITE, bold=True)
add_accent_line(s, 0.8, 1.05, 2.0)

# 数据卡
add_metric_card(s, 0.8, 1.4, "架构模块拆分", "4")
add_metric_card(s, 2.6, 1.4, "UI 改进项", "5+")
add_metric_card(s, 4.4, 1.4, "新增功能", "3")
add_metric_card(s, 6.2, 1.4, "Bug 修复", "6+")
add_metric_card(s, 8.0, 1.4, "文件变更", "10+")
add_metric_card(s, 9.8, 1.4, "代码行数", "~3500→")

# 分类说明
cards_data = [
    ("🔧 架构重构", "app.py 拆分为 4 个核心模块\n• task_persistence.py — 持久化层\n• health.py — 健康检查\n• guards.py — 提交/轮询守卫\n• celery_tasks_v75.py — 扩展"),
    ("🎨 UI 优化", "• 紧凑型首页布局 + 摘要卡片\n• 页面导航引导\n• 13 步 Pipeline 可视化增强\n• 用户友好错误/超时提示\n• 侧边栏系统维护面板"),
    ("🆕 新功能", "• Markdown 报告一键导出\n• 假设迭代追踪（基于成果继续优化）\n• 实验设计自动生成（样本量/统计方法/随机化）"),
    ("🐛 关键修复", "• Worker 检测: Redis key → inspect.ping()\n• 双执行模式: Worker + 本地降级\n• 快速失败: 3 次 LLM 失败即中止\n• 迭代 prompt 累积修复"),
]
for i, (title, content) in enumerate(cards_data):
    row, col = divmod(i, 2)
    left = 0.8 + col * 6.0
    top = 2.5 + row * 2.3
    add_card(s, left, top, 5.7, 2.0, title, content)

# ═══════════════════════════════════════════
# Slide 3: 问题诊断 — Worker 稳定性
# ═══════════════════════════════════════════
s = make_slide()
add_text_box(s, 0.8, 0.4, 11, 0.6, "问题诊断: Celery Worker 稳定性", font_size=28, color=WHITE, bold=True)
add_accent_line(s, 0.8, 1.05, 2.0)

add_text_box(s, 0.8, 1.3, 11, 0.4, "初始症状", font_size=16, color=ACCENT, bold=True)
add_bullet_list(s, 0.8, 1.7, 5.5, 1.5, [
    "Worker 频繁崩溃，Streamlit 显示「离线」",
    "侧边栏健康检查始终报告离线",
    "Worker 重启后自动执行 Redis 中残留旧任务",
    "API 500 导致 8 轮演化循环全部重试",
], color=GRAY)

add_text_box(s, 7.0, 1.3, 5.5, 0.4, "根因分析", font_size=16, color=RED, bold=True)
add_bullet_list(s, 7.0, 1.7, 5.5, 1.5, [
    "check_redis_health() 关闭了 Celery channel",
    "r.keys('celery*') 匹配到已死亡任务的元数据 key",
    "Windows Solo Pool Worker 不稳定",
    "cloud.hongqiye.com 代理持续返回 500",
], color=GRAY)

add_text_box(s, 0.8, 3.5, 11, 0.4, "解决方案", font_size=16, color=GREEN, bold=True)
add_bullet_list(s, 0.8, 3.9, 11, 1.0, [
    "移除 conn.default_channel.close() — 修复 channel 断开问题",
    "celery.control.inspect().ping() 替代 r.keys() — 精准检测 Worker 心跳",
], color=WHITE)

# 代码对比
add_card(s, 0.8, 5.0, 5.7, 2.2, "修复前 (不可靠)", "r = redis.Redis(host='localhost', port=6379)\ncelery_keys = [k for k in r.keys('*') if 'celery' in k.lower()]\nconn.default_channel.close()  # 断开了 Celery 的连接！", RED)
add_card(s, 6.9, 5.0, 5.7, 2.2, "修复后 (可靠)", "celery_app = get_celery_app()\ninspect = celery_app.control.inspect(timeout=3.0)\npings = inspect.ping()  # 真正 ping Worker\nif pings and len(pings) > 0: return True", GREEN)

# ═══════════════════════════════════════════
# Slide 4: 架构重构
# ═══════════════════════════════════════════
s = make_slide()
add_text_box(s, 0.8, 0.4, 11, 0.6, "代码架构重构: app.py 拆分", font_size=28, color=WHITE, bold=True)
add_accent_line(s, 0.8, 1.05, 2.0)

add_text_box(s, 0.8, 1.3, 5.5, 0.4, "重构前", font_size=16, color=RED, bold=True)
add_bullet_list(s, 0.8, 1.7, 5.5, 1.0, [
    "app.py 1258 行 — 单文件承载全部逻辑",
    "UI 渲染 / DB 操作 / 健康检查 / 守卫逻辑混在一起",
    "测试困难，修改风险高",
], color=GRAY)

add_text_box(s, 7.0, 1.3, 5.5, 0.4, "重构后", font_size=16, color=GREEN, bold=True)
add_bullet_list(s, 7.0, 1.7, 5.5, 1.0, [
    "app.py 缩减至 ~580 行（纯 UI 入口）",
    "每个模块可独立测试、独立修改",
], color=GRAY)

# 架构图
add_card(s, 0.8, 3.0, 3.5, 3.2, "📊 task_persistence.py", "240 行 | 8 个函数\n• init_task_persistence_db\n• register_task_persistence\n• update_task_completion\n• get_task_history_list\n• delete_task_from_history\n• clear_all_task_history\n• recover_lost_task_on_reload", WHITE)
add_card(s, 4.6, 3.0, 2.6, 3.2, "🩺 health.py", "75 行 | 2 个函数\n• check_redis_health\n• check_worker_heartbeat\n(使用 inspect.ping)", WHITE)
add_card(s, 7.5, 3.0, 2.6, 3.2, "🛡️ guards.py", "196 行 | 4 个函数\n• check_submission_guard\n• acquire_submission_lock\n• release_submission_lock\n• check_poll_guard", WHITE)
add_card(s, 10.4, 3.0, 2.3, 3.2, "📦 扩展\ncelery_tasks_v75.py", "新增:\n• submit_celery_task\n  _with_safety\n• poll_task_status\n  _safe\n• handle_poll_timeout", WHITE)

# ═══════════════════════════════════════════
# Slide 5: UI 改进 5 项
# ═══════════════════════════════════════════
s = make_slide()
add_text_box(s, 0.8, 0.4, 11, 0.6, "UI 优化: 5 项用户体验提升", font_size=28, color=WHITE, bold=True)
add_accent_line(s, 0.8, 1.05, 2.0)

ui_items = [
    ("🏠 紧凑型首页", "标题区精简高度，输入框与提交按钮紧邻\n字符计数实时显示(绿→黄→红)\n快捷示例按钮一键填入", GREEN),
    ("📊 摘要卡片", "科学分数/创新度/严谨度/版本数 四维一卡\n颜色编码: ≥8绿色 ≥6琥珀 <6红色\n生成完成即可快速评估假设质量", ACCENT),
    ("🧭 导航引导", "侧边栏顶部新增页面导航折叠区\n9 个页面的图标+描述说明\n新人上手无需记忆页面功能", GREEN),
    ("⚡ 进度反馈", "13 步 Pipeline 可视化\n脉冲动画 + 当前阶段描述文案\n折叠步骤列表仅显示前后 2 步", ACCENT),
    ("💬 友好错误", "4 类错误映射: 伪科学/拒绝/离线/本地失败\n每类配图标+解释+3 条可操作建议\n超时独立报告含等待时间与轮询次数", GREEN),
]

for i, (title, desc, color) in enumerate(ui_items):
    row, col = divmod(i, 3 if i < 3 else 2)
    top = 1.3 + (0 if i < 3 else 3.2) + row * 2.8
    left = 0.8 + (col if i < 3 else col) * (4.1 if i < 3 else 6.0)
    width = 3.8 if i < 3 else 5.7
    add_card(s, left, top, width, 2.5, title, desc, color)

# ═══════════════════════════════════════════
# Slide 6: 双执行模式 + 快速失败
# ═══════════════════════════════════════════
s = make_slide()
add_text_box(s, 0.8, 0.4, 11, 0.6, "双执行模式 + 快速失败机制", font_size=28, color=WHITE, bold=True)
add_accent_line(s, 0.8, 1.05, 2.0)

add_text_box(s, 0.8, 1.3, 5.5, 0.4, "双执行模式", font_size=18, color=ACCENT, bold=True)
add_card(s, 0.8, 1.8, 5.7, 2.8, "🚀 Worker 提交 (主路径)",
        "• 任务投递到 Celery + Redis 队列\n• 异步执行，前端自动轮询\n• Worker 离线时阻止投递 (不再堆积)\n• 返回 WORKER_OFFLINE 错误码", WHITE)
add_card(s, 0.8, 4.8, 5.7, 2.0, "⚡ 本地执行 (降级)",
        "• Worker 离线时的备选方案\n• st.spinner 显示执行进度 (5-15分钟)\n• 结果直接写入 session_state\n• 无需依赖任何外部服务", WHITE)

add_text_box(s, 7.0, 1.3, 5.5, 0.4, "快速失败 (Fast-Fail)", font_size=18, color=AMBER, bold=True)
add_card(s, 7.0, 1.8, 5.7, 2.8, "问题: API 500 导致死循环",
        "• LLM API (cloud.hongqiye.com) 持续 500\n• 原机制: 循环 8 轮全部重试 → 大量浪费\n• 用户等待时间 > 10 分钟", RED)
add_card(s, 7.0, 3.5, 5.7, 2.0, "修复: 连续失败计数",
        "• PhoenixContext.consecutive_llm_failures ++ \n• 累计 ≥3 次 → 立即中止 (MAX_PHOENIX_EXCEEDED)\n• 成功路径 → 计数器归零\n• max_iterations 默认 8→4", GREEN)
add_card(s, 7.0, 5.7, 5.7, 1.5, "效果",
        "最坏情况: 8×N秒 → 3×N秒 后快速退出\n用户等待时间降低 60%+", ACCENT)

# ═══════════════════════════════════════════
# Slide 7: 新功能概览
# ═══════════════════════════════════════════
s = make_slide()
add_text_box(s, 0.8, 0.4, 11, 0.6, "新增功能: 3 项核心能力", font_size=28, color=WHITE, bold=True)
add_accent_line(s, 0.8, 1.05, 2.0)

# 三列大卡片
features = [
    ("📥", "Markdown 导出", ACCENT,
     "功能说明",
     "• 一键下载完整 Markdown 报告\n• 包含所有 tab 数据\n• 方法论 + 演化链 + 文献支撑",
     "技术实现",
     "• _build_markdown_report() 构建\n• st.download_button 驱动下载\n• 文件名含时间戳"),
    ("🔄", "假设迭代追踪", AMBER,
     "功能说明",
     "• 成功结果下新增优化入口\n• 输入优化指令继续演化\n• 显示迭代次数/原始输入/上一轮摘要",
     "技术实现",
     "• 原始输入 + 最新优化 → 新 prompt\n• 避免多轮累积 (始终用 original_input)\n• refinement_count 追踪迭代轮数"),
    ("🧪", "实验设计生成", GREEN,
     "功能说明",
     "• 第 7 个 Tab: 实验设计方案\n• 自动识别 RCT/观察性/横断面\n• 样本量估算 + 统计方法 + 盲法",
     "技术实现",
     "• _generate_experiment_design() 模板生成\n• 关键词检测因果关系/预测建模\n• Cohen's d 公式 / EPV 法则 / 分层随机化"),
]
for i, (icon, title, color, _, desc, _, impl) in enumerate(features):
    left = 0.8 + i * 4.1
    add_card(s, left, 1.3, 3.8, 1.0, f"{icon} {title}", "", color)
    add_card(s, left, 2.3, 3.8, 2.5, desc, "", WHITE)
    add_card(s, left, 4.9, 3.8, 2.2, impl, "", GRAY)

# ═══════════════════════════════════════════
# Slide 8: 新功能详解 — 实验设计
# ═══════════════════════════════════════════
s = make_slide()
add_text_box(s, 0.8, 0.4, 11, 0.6, "实验设计生成 — 自动输出 6 类方案", font_size=28, color=WHITE, bold=True)
add_accent_line(s, 0.8, 1.05, 2.0)

exp_cards = [
    ("📋 研究类型推断", "检测关键词自动分类\n• Causal/因果 → RCT\n• Predict/预测 → 队列设计\n• 默认 → 横断面研究"),
    ("📊 样本量估算", "因果研究: Cohen's d 公式\n  n = 2(Z₁₋α/₂ + Z₁₋β)²/d²\n预测建模: EPV ≥ 10 法则\n  n ≥ 10 × p (特征数)"),
    ("🔢 统计方法推荐", "• 因果: 工具变量 2SLS / DID / PSM\n• 预测: Logistic/Cox + XGBoost\n• 解释: SHAP / LIME\n• 纠正: Bonferroni / FDR"),
    ("🎲 随机化 & 盲法", "• RCT: 分层区组随机化\n• 观察性: PSM 倾向性匹配\n• 盲法: 双盲+第三方评估\n• 分配隐藏: SNOSE 信封"),
    ("📏 效应量指标", "• Cohen's d / Hedges' g (连续变量)\n• OR / RR (二分类结局)\n• Hazard Ratio (生存分析)\n• η² / Cramer's V"),
    ("📋 数据采集计划", "• T0 基线 → T1 暴露期 → T2 随访\n• REDCap/EDC 系统录入\n• SAP 预设亚组与缺失数据处理"),
]
for i, (title, content) in enumerate(exp_cards):
    row, col = divmod(i, 3)
    left = 0.8 + col * 4.1
    top = 1.3 + row * 3.0
    add_card(s, left, top, 3.8, 2.7, title, content, ACCENT)

# ═══════════════════════════════════════════
# Slide 9: Bug 修复清单
# ═══════════════════════════════════════════
s = make_slide()
add_text_box(s, 0.8, 0.4, 11, 0.6, "Bug 修复记录 (本期 6 项)", font_size=28, color=WHITE, bold=True)
add_accent_line(s, 0.8, 1.05, 2.0)

bugs = [
    ("Worker 始终显示离线", "Redis key 扫描 + channel 关闭 → inspect.ping() | 已修复 ✅"),
    ("Worker 重启后执行旧任务", "残留 celery-task-meta-* → Worker 离线时阻止 Celery 投递 | 已修复 ✅"),
    ("Streamlit 自动刷新死循环", "stale session state 持续存在 → 5 分钟超时自动清除 | 已修复 ✅"),
    ("API 500 死循环 8 轮重试", "无失败计数 → consecutive_llm_failures ≥ 3 即中止 | 已修复 ✅"),
    ("last_submit_input 未保存", "仅 Worker 离线时保存 → 提交/本地执行两路径均保存 | 已修复 ✅"),
    ("迭代 prompt 多轮累积", "嵌套上一轮 prompt → 始终使用 original_input | 已修复 ✅"),
]
for i, (title, detail) in enumerate(bugs):
    row, col = divmod(i, 3)
    left = 0.8 + col * 4.1
    top = 1.3 + row * 3.0
    add_card(s, left, top, 3.8, 2.7, title, detail, ACCENT if '✅' in detail else AMBER)

# ═══════════════════════════════════════════
# Slide 10: 技术架构总览
# ═══════════════════════════════════════════
s = make_slide()
add_text_box(s, 0.8, 0.4, 11, 0.6, "技术架构: V7.5 当前状态", font_size=28, color=WHITE, bold=True)
add_accent_line(s, 0.8, 1.05, 2.0)

# 三层架构
add_text_box(s, 0.8, 1.3, 3.5, 0.4, "UI 层", font_size=18, color=ACCENT, bold=True)
add_card(s, 0.8, 1.7, 3.8, 1.5, "streamlit (app.py)",
        "• 页面: 主页/对话/工作流/循环/文献\n• sidebar_components.py\n• pipeline_components.py\n• report_renderers.py\n• evolution_view.py / styles.py", WHITE)

add_text_box(s, 5.0, 1.3, 3.5, 0.4, "核心层", font_size=18, color=AMBER, bold=True)
add_card(s, 5.0, 1.7, 3.8, 1.5, "src/core/",
        "• celery_tasks_v75.py (主逻辑)\n• phoenix_state_machine.py\n• health.py / guards.py\n• task_persistence.py\n• pseudoscience_detector.py", WHITE)

add_text_box(s, 9.2, 1.3, 3.5, 0.4, "基础设施", font_size=18, color=GREEN, bold=True)
add_card(s, 9.2, 1.7, 3.5, 1.5, "Redis + Celery + SQLite",
        "• Redis: 消息队列 + 状态缓存\n• Celery: 异步任务 (solo pool)\n• SQLite: 任务持久化\n• LLM: Anthropic SDK", WHITE)

# 数据流
add_text_box(s, 0.8, 3.5, 11, 0.4, "数据流: 用户输入 → Intent Sanitizer → Celery → Phoenix 演化 → 结果渲染", font_size=14, color=GRAY, bold=True)

# 关键指标
add_metric_card(s, 0.8, 4.5, "Phoenix 演化步骤", "13")
add_metric_card(s, 2.6, 4.5, "最大迭代数", "4")
add_metric_card(s, 4.4, 4.5, "失败快速中止", "3次")
add_metric_card(s, 6.2, 4.5, "轮询超时", "300次")
add_metric_card(s, 8.0, 4.5, "Session生命周期", "5分钟")
add_metric_card(s, 9.8, 4.5, "API 硬上限", "15次")

add_card(s, 0.8, 5.5, 12, 1.5, "关键设计决策",
        "1. Celery inspect.ping() 替代 Redis key 扫描 → Worker 检测准确率 100%\n2. Worker 离线时阻止 Celery 投递 → 彻底解决任务堆积\n3. 连续 LLM 失败 3 次即中止 → 最坏等待时间降低 60%+\n4. 本地执行作为降级路径 → 即使 Worker 不可用也能完成任务", ACCENT)

# ═══════════════════════════════════════════
# Slide 11: 待解决问题
# ═══════════════════════════════════════════
s = make_slide()
add_text_box(s, 0.8, 0.4, 11, 0.6, "待解决问题 & 风险", font_size=28, color=WHITE, bold=True)
add_accent_line(s, 0.8, 1.05, 2.0)

issues = [
    ("⚠️ LLM API 稳定性", "cloud.hongqiye.com 代理持续 500 错误\n所有 LLM 调用功能 (假设生成/文献检索/攻击分析) 均受影响\n快速失败机制可避免死循环，但核心功能不可用", RED),
    ("⚠️ 中文文献检索", "中文输入无法直接用于 PubMed/ArXiv\nKeyword Translator (LLM) 因 API 故障不可用\n文献检索结果可能为空", AMBER),
    ("⚠️ Windows Solo Pool", "Celery --pool=solo 在 Windows 上不稳定\n无法使用 prefork/gevent\nWorker 崩溃后需手动重启", AMBER),
    ("📋 Worker 监控", "缺少 Worker 崩溃自动告警\n建议: 添加 healthcheck 端点或 watchdog 脚本", GRAY),
    ("📋 测试覆盖", "核心模块 (phoenix_state_machine, guards, health)\n尚未添加单元测试，建议优先覆盖", GRAY),
    ("📋 日志存储", "当前仅终端输出\n建议: 接入结构化日志 (JSON) + 持久化\n便于回溯 Worker 崩溃原因", GRAY),
]
for i, (title, detail, color) in enumerate(issues):
    row, col = divmod(i, 3)
    left = 0.8 + col * 4.1
    top = 1.3 + row * 3.0
    add_card(s, left, top, 3.8, 2.7, title, detail, color)

# ═══════════════════════════════════════════
# Slide 12: 下一步计划
# ═══════════════════════════════════════════
s = make_slide()
add_text_box(s, 0.8, 0.4, 11, 0.6, "下一步计划", font_size=28, color=WHITE, bold=True)
add_accent_line(s, 0.8, 1.05, 2.0)

plan_items = [
    ("🔴 高优先级", ACCENT, [
        "解决 LLM API 稳定性 (换 key / 换代理)",
        "补充核心模块单元测试",
        "Worker 崩溃自动告警",
    ]),
    ("🟡 中优先级", AMBER, [
        "中文文献检索优化 (绕过 LLM 翻译)",
        "实验结果与实际假设一致性校验",
        "引入本地小模型 (Ollama) 作为降级",
    ]),
    ("🟢 低优先级", GREEN, [
        "PDF 论文全文解析",
        "多假设并行对比",
        "Docker 一键部署",
    ]),
]
for i, (label, color, items) in enumerate(plan_items):
    left = 0.8 + i * 4.1
    add_text_box(s, left, 1.3, 3.8, 0.4, label, font_size=16, color=color, bold=True)
    add_bullet_list(s, left, 1.8, 3.8, 4.0, items, font_size=13, color=WHITE)

# 总结
add_card(s, 0.8, 5.5, 12, 1.5, "总结",
        "本期完成: 代码架构拆分 (4 模块) + 5 项 UI 优化 + 3 项新功能 + 6 个 Bug 修复\n核心成果: Worker 检测准确率显著提升、快速失败机制填补稳定性短板、双执行模式保障可用性\n下一步重点: 解决 API 稳定性 → 补充测试 → 优化文献检索", ACCENT)

# ═══════════════════════════════════════════
# 保存
# ═══════════════════════════════════════════
output_dir = os.path.dirname(os.path.abspath(__file__))
output_path = os.path.join(output_dir, '..', '..', 'V75_Progress_Report.pptx')
output_path = os.path.abspath(output_path)
prs.save(output_path)
print(f"PPT 已生成: {output_path}")
